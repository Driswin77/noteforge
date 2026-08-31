import sys
import json
import os

def extract_pdf(filepath):
    import pdfplumber
    slides = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            slides.append({
                "slide_number": i,
                "text": text.strip()
            })
    return slides

def extract_pptx(filepath):
    from pptx import Presentation
    prs = Presentation(filepath)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        texts.append(line.strip())
        slides.append({
            "slide_number": i,
            "text": "\n".join(texts)
        })
    return slides

def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "No file path provided"}))
        sys.exit(1)

    filepath = sys.argv[1]

    if not os.path.exists(filepath):
        print(json.dumps({"error": f"File not found: {filepath}"}))
        sys.exit(1)

    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == ".pdf":
            slides = extract_pdf(filepath)
        elif ext == ".pptx":
            slides = extract_pptx(filepath)
        else:
            print(json.dumps({"error": f"Unsupported file type: {ext}"}))
            sys.exit(1)

        result = {
            "filename": os.path.basename(filepath),
            "total_slides": len(slides),
            "slides": slides
        }
        print(json.dumps(result, indent=2))

    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)

if __name__ == "__main__":
    main()